#!/usr/bin/env python3
"""Convert a compiled Beamer PDF into a PowerPoint deck.

Each PDF page is rasterized to PNG at high DPI and placed as a full-bleed image
on a 16:9 PPTX slide. Output is visually identical to the PDF but the slides
are not text-editable (they are image slides).

Usage:
    python scripts/pdf_to_pptx.py writeup/slides.pdf writeup/slides.pptx
    python scripts/pdf_to_pptx.py writeup/slides.pdf writeup/slides.pptx --dpi 200

Requires:
    pip install python-pptx
    brew install poppler          # provides pdftoppm
"""
from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches


SLIDE_W_IN = 13.3333  # 16:9 widescreen, matches Beamer aspectratio=169
SLIDE_H_IN = 7.5


def rasterize_pdf(pdf_path: Path, out_dir: Path, dpi: int) -> list[Path]:
    """Call pdftoppm to render every page of the PDF as PNG."""
    if shutil.which("pdftoppm") is None:
        sys.exit("error: pdftoppm not found. Install poppler (brew install poppler).")
    prefix = out_dir / "page"
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(prefix)],
        check=True,
    )
    return sorted(out_dir.glob("page-*.png"))


def build_pptx(images: list[Path], output: Path) -> None:
    """Create a PPTX where each slide is a full-bleed PNG from the PDF."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]  # fully blank layout

    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output))


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pdf", type=Path, help="Compiled Beamer PDF")
    parser.add_argument("pptx", type=Path, help="Output PPTX path")
    parser.add_argument("--dpi", type=int, default=200, help="Raster DPI (default 200)")
    args = parser.parse_args()

    if not args.pdf.is_file():
        sys.exit(f"error: {args.pdf} not found")
    args.pptx.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        print(f"[1/2] rasterising {args.pdf.name} @ {args.dpi} dpi ...")
        pages = rasterize_pdf(args.pdf, tmp_dir, args.dpi)
        if not pages:
            sys.exit("error: pdftoppm produced no pages")
        print(f"      {len(pages)} pages rendered")

        print(f"[2/2] writing {args.pptx} ...")
        build_pptx(pages, args.pptx)

    mb = args.pptx.stat().st_size / 1024 / 1024
    print(f"done. {args.pptx} ({mb:.1f} MB)")


if __name__ == "__main__":
    main()
